import matplotlib.pyplot as plt
import networkx as nx
from pptx import Presentation
from pptx.util import Inches, Pt
import os

# ---------------------------------------------------------------
# BUILD EDGES FROM ANY NESTED JSON-LIKE DICTIONARY
# ---------------------------------------------------------------
def build_edges(data, parent=None, edges=None):
    if edges is None:
        edges = []

    if isinstance(data, dict):
        for k, v in data.items():
            if parent:
                edges.append((parent, k))
            build_edges(v, k, edges)
    elif isinstance(data, list):
        for item in data:
            if parent:
                edges.append((parent, item))
            build_edges(item, item, edges)

    return edges


# ---------------------------------------------------------------
# CHOOSE LAYOUT
# ---------------------------------------------------------------
def get_layout(G, layout_type):
    if layout_type == "flowchart":
        return nx.nx_agraph.graphviz_layout(G, prog="dot")

    if layout_type == "tree":
        return nx.nx_agraph.graphviz_layout(G, prog="twopi")

    if layout_type == "mindmap":
        pos = nx.nx_agraph.graphviz_layout(G, prog="dot")
        return {k: (x * 1.5, -y) for k, (x, y) in pos.items()}

    if layout_type == "circular":
        return nx.circular_layout(G)

    raise ValueError(f"Unknown layout: {layout_type}")


# ---------------------------------------------------------------
# EXPORTS: PNG / SVG / PDF
# ---------------------------------------------------------------
def save_matplotlib(G, pos, file_prefix):
    plt.figure(figsize=(12, 8))

    nx.draw(
        G, pos,
        with_labels=True,
        arrows=True,
        node_size=3000,
        node_color="white",
        edge_color="black",
        linewidths=1.5,
        font_size=9
    )

    plt.tight_layout()
    plt.savefig(f"{file_prefix}.png", dpi=300)
    plt.savefig(f"{file_prefix}.svg")
    plt.savefig(f"{file_prefix}.pdf")
    plt.close()


# ---------------------------------------------------------------
# OPTIONAL: Generate editable PPTX (enable when needed)
# ---------------------------------------------------------------
def save_pptx(G, pos, file_prefix):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    min_x = min(p[0] for p in pos.values())
    max_x = max(p[0] for p in pos.values())
    min_y = min(p[1] for p in pos.values())
    max_y = max(p[1] for p in pos.values())

    def norm(v, vmin, vmax):
        return (v - vmin) / (vmax - vmin + 1e-6)

    node_positions = {}

    # Add nodes
    for node, (x, y) in pos.items():
        left = Inches(norm(x, min_x, max_x) * 9)
        top = Inches(norm(-y, min_y, max_y) * 5)
        shape = slide.shapes.add_shape(1, left, top, Inches(1.5), Inches(0.6))
        shape.text = node
        node_positions[node] = (left, top)

    # Add connectors
    for (u, v) in G.edges():
        start_x, start_y = node_positions[u]
        end_x, end_y = node_positions[v]
        line = slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
        line.line.width = Pt(1.2)

    prs.save(f"{file_prefix}.pptx")


# ---------------------------------------------------------------
# MAIN FUNCTION: YOU CALL ONLY THIS
# ---------------------------------------------------------------
def generate_diagram(api_response, layout_type, generate_pptx=False):
    edges = build_edges(api_response)
    G = nx.DiGraph()
    G.add_edges_from(edges)

    pos = get_layout(G, layout_type)
    file_prefix = f"diagram_{layout_type}"

    # export png/svg/pdf
    save_matplotlib(G, pos, file_prefix)

    # ppt OPTIONAL
    if generate_pptx:
        save_pptx(G, pos, file_prefix)

    return f"Generated {file_prefix}.*"


# ---------------------------------------------------------------
# EXAMPLE (your API response)
# ---------------------------------------------------------------
if _name_ == "_main_":
    api_response = {
        "Home": {
            "Login": {
                "Validate": ["Success", "Failure"]
            },
            "Menu": ["Dashboard", "Reports", "Settings"]
        }
    }

    # Choose ONE:
    layout = "flowchart"      # flowchart | tree | mindmap | circular

    generate_diagram(api_response, layout, generate_pptx=False)